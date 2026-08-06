"""Markdown -> .pptx rendering (PHASE-3-SPEC.md Part 5).

Decks are rendered from the stored markdown, never from a second query. That is the whole
design: a generator decides what is in a document once, under the promotion and stamping rules
(app/generators.py), and this module only lays it out. A renderer that re-queried could quietly
include a record the generator excluded, which is exactly the class of bug the "by construction"
rule exists to prevent.

Binaries are produced on download and never stored — the markdown is the artifact.
"""
from __future__ import annotations

import io
import re
import textwrap

from pptx import Presentation
from pptx.util import Inches, Pt

# Deliberately plain: no theme, no colour. These decks are drafts the operator adapts, and a
# half-branded template reads worse than an unstyled one. Status colour would also import the
# design system's semantics into a surface that has none of its guardrails.
_TITLE_SIZE = Pt(30)
_BODY_SIZE = Pt(16)
_META_SIZE = Pt(11)
_BULLET_LIMIT = 9          # past this a slide stops being readable; the rest spills to a continuation


def _blocks(markdown: str) -> list[tuple[str, list[str]]]:
    """Split markdown into (heading, lines) blocks on '## '. Content before the first heading
    belongs to the title slide."""
    blocks: list[tuple[str, list[str]]] = []
    heading, buf = None, []
    for line in markdown.splitlines():
        if line.startswith("## "):
            blocks.append((heading, buf))
            heading, buf = line[3:].strip(), []
        else:
            buf.append(line)
    blocks.append((heading, buf))
    return blocks


def _bullets(lines: list[str]) -> list[str]:
    """Turn a markdown block into speakable bullets.

    Tables become one bullet per row with the header names inlined, because a pptx table that
    nobody sized is worse than a sentence. Emphasis and pipes are stripped; the text carries.
    """
    out: list[str] = []
    headers: list[str] = []
    for raw in lines:
        line = raw.strip()
        if not line or line.startswith(">"):
            if line.startswith(">"):
                out.append(line.lstrip("> ").strip())
            continue
        if re.fullmatch(r"\|[\s|:-]+\|", line):        # table separator row
            continue
        if line.startswith("|"):
            cells = [c.strip() for c in line.strip("|").split("|")]
            if not headers:
                headers = cells
                continue
            pairs = [f"{h}: {c}" for h, c in zip(headers, cells) if c and c != "—"]
            if pairs:
                out.append(" · ".join(pairs))
            continue
        headers = []
        if line.startswith(("- ", "* ")):
            out.append(line[2:].strip())
        elif line.startswith("_") and line.endswith("_"):
            out.append(line.strip("_"))
        else:
            out.append(line)
    # Strip emphasis markers wherever they appear, not just as whole-line wrappers: a line like
    # "_Recommendation._ 10,706 seats remain" is emphasis mid-sentence and used to reach the
    # slide with the underscores intact.
    cleaned = []
    for b in out:
        b = re.sub(r"\*\*(.+?)\*\*", r"\1", b)
        b = re.sub(r"(?<![A-Za-z0-9])_(.+?)_(?![A-Za-z0-9])", r"\1", b)
        b = re.sub(r"[*`]", "", b).strip()
        if b and b != "None.":
            cleaned.append(b)
    return cleaned


def render(markdown: str, *, title: str, subtitle: str = "") -> bytes:
    """Render stored markdown to a .pptx byte stream."""
    prs = Presentation()
    blank = prs.slide_layouts[6]

    def slide_with(heading: str, bullets: list[str]):
        s = prs.slides.add_slide(blank)
        box = s.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(1))
        tf = box.text_frame
        tf.text = heading
        tf.paragraphs[0].runs[0].font.size = _TITLE_SIZE
        tf.paragraphs[0].runs[0].font.bold = True
        body = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(9), Inches(5))
        btf = body.text_frame
        btf.word_wrap = True
        for i, b in enumerate(bullets):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = _BODY_SIZE
        return s

    # Title slide, carrying the stamp — a deck without its generation date is a deck that gets
    # presented six weeks stale.
    first = prs.slides.add_slide(blank)
    tb = first.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.5), Inches(2))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].runs[0].font.size = Pt(40)
    tf.paragraphs[0].runs[0].font.bold = True
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = _META_SIZE

    for heading, lines in _blocks(markdown):
        if heading is None:
            continue
        bullets = _bullets(lines) or ["None on record."]
        for i in range(0, len(bullets), _BULLET_LIMIT):
            chunk = bullets[i:i + _BULLET_LIMIT]
            slide_with(heading if i == 0 else f"{heading} (cont.)", chunk)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_pdf(markdown: str, *, title: str, subtitle: str = "") -> bytes:
    """Render the stored artifact to a dependency-free, text-first PDF.

    The markdown remains canonical. This intentionally modest renderer preserves every claim and
    citation and produces a valid, selectable-text PDF without introducing a browser/HTML engine.
    """
    raw_lines = [title, subtitle, ""]
    for heading, lines in _blocks(markdown):
        if heading:
            raw_lines += [heading.upper(), ""]
        raw_lines += _bullets(lines)
        raw_lines.append("")

    lines: list[str] = []
    for raw in raw_lines:
        if not raw:
            lines.append("")
            continue
        lines.extend(textwrap.wrap(str(raw), width=92, break_long_words=False,
                                   break_on_hyphens=False) or [""])
    pages = [lines[i:i + 48] for i in range(0, len(lines), 48)] or [[]]

    def pdf_text(value: str) -> str:
        value = value.replace("—", "-").replace("→", "->").replace("·", "-")
        value = value.encode("latin-1", "replace").decode("latin-1")
        return value.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")

    objects: dict[int, bytes] = {
        1: b"<< /Type /Catalog /Pages 2 0 R >>",
        3: b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    }
    kids = []
    for index, page_lines in enumerate(pages):
        page_id = 4 + index * 2
        content_id = page_id + 1
        kids.append(f"{page_id} 0 R")
        commands = ["BT", "/F1 10 Tf", "14 TL", "50 755 Td"]
        for line in page_lines:
            commands.append(f"({pdf_text(line)}) Tj")
            commands.append("T*")
        commands.append("ET")
        stream = "\n".join(commands).encode("latin-1")
        objects[page_id] = (f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
                            f"/Resources << /Font << /F1 3 0 R >> >> /Contents {content_id} 0 R >>").encode()
        objects[content_id] = (f"<< /Length {len(stream)} >>\nstream\n".encode() + stream +
                               b"\nendstream")
    objects[2] = f"<< /Type /Pages /Kids [{' '.join(kids)}] /Count {len(kids)} >>".encode()

    out = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets = [0]
    for oid in range(1, max(objects) + 1):
        offsets.append(len(out))
        out += f"{oid} 0 obj\n".encode() + objects[oid] + b"\nendobj\n"
    xref = len(out)
    out += f"xref\n0 {len(offsets)}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets[1:]:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (f"trailer\n<< /Size {len(offsets)} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n").encode()
    return bytes(out)
