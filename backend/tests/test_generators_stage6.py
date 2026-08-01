"""Stage 6 — generators as finished artifacts (PHASE-3-SPEC.md Part 5, EXPANSION-ENGINE-SPEC §8).

The tests that matter are the ones guarding what may leave the building: promotion enforced in
the query rather than by review, stamping on every artifact, and drafts that stay drafts.
"""
import os
import tempfile
from datetime import date, timedelta

import pytest
from fastapi.testclient import TestClient


@pytest.fixture()
def client():
    fd, path = tempfile.mkstemp(suffix=".sqlite")
    os.close(fd)
    os.environ["VALENCE_OS_DB"] = path
    os.environ["VALENCE_OS_WORKER"] = "0"
    from app.main import app
    with TestClient(app) as c:
        yield c
    for s in ("", "-wal", "-shm"):
        try: os.unlink(path + s)
        except FileNotFoundError: pass


def _today():
    from app.db import now_utc
    return now_utc()[:10]


def _days(n):
    return (date.fromisoformat(_today()) + timedelta(days=n)).isoformat()


@pytest.fixture()
def acct(client):
    a = client.post("/api/accounts", json={"name": "Terravance"}).json()
    p = client.post("/api/programs", json={"account_id": a["id"], "name": "Global"}).json()
    part = client.post("/api/population-partitions", json={
        "account_id": a["id"], "total_fte": 20000}).json()
    seg = client.post("/api/population-segments", json={
        "partition_id": part["id"], "name": "DACH", "headcount": 6000}).json()
    seg2 = client.post("/api/population-segments", json={
        "partition_id": part["id"], "name": "Nordics", "headcount": 4000}).json()
    uc = client.post("/api/use-cases", json={"name": "Performance reviews", "slug": "pr"}).json()
    person = client.post("/api/persons", json={
        "name": "Dana Okafor", "affiliation": "client", "account_id": a["id"]}).json()
    client.post("/api/stakeholder-roles", json={
        "program_id": p["id"], "person_id": person["id"], "role": "champion",
        "stance": "supporter", "stance_assessed_on": _today(), "stance_evidence_note": "e"})
    source = client.post("/api/source-references", json={
        "label": "Mock scorecard", "type": "data_report"}).json()
    return {"a": a, "p": p, "seg": seg, "seg2": seg2, "uc": uc,
            "person": person, "source": source}


# --- promotion, enforced in the query ----------------------------------------------------------
def test_business_case_carries_only_promoted_evidence(client, acct):
    aid = acct["a"]["id"]
    client.post("/api/value-stories", json={
        "account_id": aid, "outcome": "PROMOTED win", "visibility_class": "qbr_exec",
        "evidence_tier": "measured_operational", "source_reference_id": acct["source"]["id"]})
    client.post("/api/value-stories", json={
        "account_id": aid, "outcome": "INTERNAL note", "visibility_class": "internal"})
    client.post("/api/value-stories", json={
        "account_id": aid, "outcome": "NEGATIVE objection", "visibility_class": "qbr_exec",
        "is_negative": True})

    bc = client.get(f"/api/accounts/{aid}/business-case").json()
    outcomes = [e["outcome"] for e in bc["evidence"]]
    assert "PROMOTED win" in outcomes
    assert "INTERNAL note" not in outcomes
    assert "NEGATIVE objection" not in outcomes
    assert bc["audience"] == "client_facing"
    assert "PROMOTED win" in bc["markdown"] and "INTERNAL note" not in bc["markdown"]


def test_business_case_scorecard_uses_only_client_accepted_targets(client, acct):
    """An internal target is not a promise the client made, so it cannot appear in a document
    arguing they got what was promised."""
    aid = acct["a"]["id"]
    d = client.post("/api/metric-definitions", json={"name": "Activation"}).json()
    client.post("/api/value-targets", json={
        "account_id": aid, "definition_id": d["id"], "segment_id": acct["seg"]["id"],
            "target_value": 0.7, "timeframe_end": _days(30),
            "client_accepted": True, "accepted_by_person_id": acct["person"]["id"],
            "accepted_on": _today(), "client_visible": True,
            "source_reference_id": acct["source"]["id"]})
    client.post("/api/value-targets", json={
        "account_id": aid, "definition_id": d["id"], "segment_id": acct["seg2"]["id"],
        "target_value": 0.9, "timeframe_end": _days(30)})          # internal, never accepted

    bc = client.get(f"/api/accounts/{aid}/business-case").json()
    assert len(bc["scorecard"]) == 1
    assert bc["scorecard"][0]["population"] == "DACH"


def test_champion_kit_is_stricter_than_the_qbr(client, acct):
    """A champion presents this without us in the room, so qbr_exec is not good enough."""
    aid = acct["a"]["id"]
    client.post("/api/value-stories", json={
        "account_id": aid, "outcome": "QBR ONLY story", "visibility_class": "qbr_exec",
        "evidence_tier": "measured_operational", "source_reference_id": acct["source"]["id"]})
    kit = client.get(f"/api/accounts/{aid}/champion-kit").json()
    assert kit["value_summary"] == []
    assert any("externally-referenceable" in m for m in kit["stamp"]["missing_or_stale_sources"])

    client.post("/api/value-stories", json={
        "account_id": aid, "outcome": "SAFE TO SHARE story",
        "visibility_class": "externally_referenceable", "evidence_tier": "measured_operational",
        "source_reference_id": acct["source"]["id"]})
    kit = client.get(f"/api/accounts/{aid}/champion-kit").json()
    assert [v["outcome"] for v in kit["value_summary"]] == ["SAFE TO SHARE story"]


def test_pre_call_brief_is_internal_and_carries_judgments(client, acct):
    """The brief is the one artifact that SHOULD carry stance and raw judgment — it never
    leaves the building, and it is useless without them."""
    b = client.get(f"/api/accounts/{acct['a']['id']}/pre-call-brief").json()
    assert b["audience"] == "internal"
    dana = next(a for a in b["attendees"] if a["name"] == "Dana Okafor")
    assert dana["stance"] == "supporter"
    assert dana["type"] == "internal_interpretation"        # labeled, not asserted as fact
    assert dana["stance_assessed_on"] == _today()           # a judgment always carries its date


# --- the expansion frame (EXPANSION-ENGINE-SPEC §8) ---------------------------------------------
def test_named_lines_exclude_already_paid_cells(client, acct):
    """A paid-but-unevidenced cell is not an expansion line. Listing it would pad the ask with
    seats the client already bought while implying value we have not demonstrated."""
    aid = acct["a"]["id"]
    paid = client.post("/api/whitespace-cells", json={
        "account_id": aid, "segment_id": acct["seg"]["id"], "use_case_id": acct["uc"]["id"],
        "paid_seats": 900}).json()
    client.post(f"/api/whitespace-cells/{paid['id']}/set-fact", json={
        "fact": "penetration", "value": "paid", "reason": "signed"})
    proven = client.post("/api/whitespace-cells", json={
        "account_id": aid, "segment_id": acct["seg2"]["id"], "use_case_id": acct["uc"]["id"],
        "client_visible": True, "source_reference_id": acct["source"]["id"]}).json()
    client.post(f"/api/whitespace-cells/{proven['id']}/set-fact", json={
        "fact": "evidence_state", "value": "anecdotal", "reason": "pilot readout"})

    bc = client.get(f"/api/accounts/{aid}/business-case").json()
    populations = {l["population"] for l in bc["lines"]}
    assert populations == {"Nordics"}          # not DACH, which is already paid


def test_expansion_frame_names_why_a_cell_is_adjacent_to_proof(client, acct):
    """"Value achieved here, projected there" is only arguable when the adjacency is stated."""
    aid = acct["a"]["id"]
    d = client.post("/api/metric-definitions", json={"name": "Activation", "stale_after_days": 30}).json()
    client.post("/api/metric-observations", json={
        "definition_id": d["id"], "program_id": acct["p"]["id"], "value": 0.9,
        "current_through": _days(-1), "population_segment_id": acct["seg"]["id"],
        "source_reference_id": acct["source"]["id"]})
    client.post("/api/value-targets", json={
        "account_id": aid, "definition_id": d["id"], "segment_id": acct["seg"]["id"],
        "target_value": 0.7, "timeframe_end": _days(30), "client_accepted": True,
        "accepted_by_person_id": acct["person"]["id"], "accepted_on": _today(),
        "client_visible": True, "source_reference_id": acct["source"]["id"]})
    # DACH is penetrated on this use case; Nordics is a proven cell on the same use case.
    dach = client.post("/api/whitespace-cells", json={
        "account_id": aid, "segment_id": acct["seg"]["id"], "use_case_id": acct["uc"]["id"],
        "paid_seats": 900}).json()
    client.post(f"/api/whitespace-cells/{dach['id']}/set-fact", json={
        "fact": "penetration", "value": "paid", "reason": "signed"})
    client.post(f"/api/whitespace-cells/{dach['id']}/set-fact", json={
        "fact": "evidence_state", "value": "measured", "reason": "readout"})
    nordics = client.post("/api/whitespace-cells", json={
        "account_id": aid, "segment_id": acct["seg2"]["id"], "use_case_id": acct["uc"]["id"],
        "client_visible": True, "source_reference_id": acct["source"]["id"]}).json()
    client.post(f"/api/whitespace-cells/{nordics['id']}/set-fact", json={
        "fact": "evidence_state", "value": "anecdotal", "reason": "pilot"})

    vr = client.get(f"/api/accounts/{aid}/value-review").json()
    frame = next(f for f in vr["expansion_frame"] if f["population"] == "Nordics")
    assert frame["adjacent_to_proof"] is True
    assert "already proven elsewhere" in frame["basis"]


def test_value_review_flags_a_missing_economic_layer(client, acct):
    vr = client.get(f"/api/accounts/{acct['a']['id']}/value-review").json()
    assert vr["attendance"]["economic_layer_covered"] is False
    assert "Economic-layer" in vr["attendance"]["flag"]

    client.post("/api/stakeholder-roles", json={
        "program_id": acct["p"]["id"], "person_id": acct["person"]["id"],
        "role": "budget_owner", "layer": "economic"})
    vr = client.get(f"/api/accounts/{acct['a']['id']}/value-review").json()
    assert vr["attendance"]["economic_layer_covered"] is True
    assert vr["attendance"]["flag"] is None


# --- documents: drafts stay drafts ---------------------------------------------------------------
def test_generated_document_saves_as_a_draft_with_its_stamp(client, acct):
    aid = acct["a"]["id"]
    doc = client.post(f"/api/accounts/{aid}/documents", json={"kind": "business_case"}).json()
    assert doc["status"] == "draft"
    assert doc["generated_at"] and doc["data_current_through"] is None
    assert doc["missing_or_stale_note"]
    assert doc["audience"] == "client_facing"
    assert client.get(f"/api/documents?account_id={aid}&status=draft").json()


def test_reviewing_a_document_records_who_and_when(client, acct):
    aid = acct["a"]["id"]
    doc = client.post(f"/api/accounts/{aid}/documents", json={"kind": "value_review"}).json()
    assert client.post(f"/api/documents/{doc['id']}/status",
                       json={"status": "reviewed"}).status_code == 422   # who reviewed it?
    ok = client.post(f"/api/documents/{doc['id']}/status",
                     json={"status": "reviewed", "reviewed_by": "operator"}).json()
    assert ok["status"] == "reviewed" and ok["reviewed_by"] == "operator" and ok["reviewed_on"]


def test_scheduled_team_update_lands_as_a_draft_never_sent(client, acct):
    """The whole reason this runs through the job table: a timer produces a draft, and a human
    decides whether it goes anywhere."""
    from app import generators, jobs
    conn = client.app.state.conn
    job = generators.schedule_weekly_update(conn)
    jobs.run_pending(conn)

    assert jobs.get_job(conn, job["id"])["status"] == "succeeded"
    drafts = client.get("/api/documents?status=draft").json()
    update = next(d for d in drafts if d["kind"] == "team_update")
    assert update["status"] == "draft" and update["audience"] == "internal"
    assert update["account_id"] is None          # portfolio-wide
    followups = [j for j in jobs.list_jobs(conn, status="queued")
                 if j["kind"] == "weekly_team_update"]
    assert len(followups) == 1 and followups[0]["scheduled_for"]


# --- pptx rendering ------------------------------------------------------------------------------
def test_document_renders_to_a_real_pptx(client, acct):
    aid = acct["a"]["id"]
    client.post("/api/value-stories", json={
        "account_id": aid, "outcome": "PROMOTED win", "visibility_class": "qbr_exec",
        "evidence_tier": "measured_operational", "source_reference_id": acct["source"]["id"]})
    doc = client.post(f"/api/accounts/{aid}/documents", json={"kind": "business_case"}).json()
    r = client.get(f"/api/documents/{doc['id']}/pptx")
    assert r.status_code == 200
    assert r.headers["content-type"].startswith("application/vnd.openxmlformats")
    assert r.content[:2] == b"PK"                     # a real zip, i.e. a real pptx

    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(r.content))
    text = "\n".join(sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame)
    assert "PROMOTED win" in text
    assert "current through" in text                  # the stamp travels with the deck


def test_deck_strips_markdown_emphasis_from_slides(client, acct):
    """A slide that reads "_Recommendation._" has leaked its source format."""
    from app import decks
    import io
    from pptx import Presentation
    md = "# T\n\n## Section\n\n_Recommendation._ Something **bold** here.\n"
    prs = Presentation(io.BytesIO(decks.render(md, title="T")))
    text = "\n".join(sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame)
    assert "Recommendation. Something bold here." in text
    assert "_" not in text and "**" not in text


def test_kickoff_deck_downloads_as_pptx(client, acct):
    r = client.get(f"/api/accounts/{acct['a']['id']}/kickoff-deck/pptx")
    assert r.status_code == 200 and r.content[:2] == b"PK"


# --- adversarial boundaries and finished workflow ----------------------------------------------
def test_generators_reject_cross_account_people_and_programs(client, acct):
    other = client.post("/api/accounts", json={"name": "Other Co"}).json()
    other_program = client.post("/api/programs", json={
        "account_id": other["id"], "name": "Other program"}).json()
    other_person = client.post("/api/persons", json={
        "account_id": other["id"], "name": "Other Person", "affiliation": "client"}).json()
    aid = acct["a"]["id"]

    assert client.get(f"/api/accounts/{aid}/pre-call-brief?program_id={other_program['id']}").status_code == 422
    assert client.get(f"/api/accounts/{aid}/pre-call-brief?person_ids={other_person['id']}").status_code == 422
    assert client.get(f"/api/accounts/{aid}/kickoff-deck?program_id={other_program['id']}").status_code == 422
    assert client.post(f"/api/accounts/{aid}/documents", json={
        "kind": "kickoff_deck", "program_id": other_program["id"]}).status_code == 422


def test_program_scoped_brief_does_not_mix_program_risks(client, acct):
    aid = acct["a"]["id"]
    sibling = client.post("/api/programs", json={"account_id": aid, "name": "Sibling"}).json()
    client.post("/api/risks", json={"program_id": acct["p"]["id"], "description": "RIGHT PROGRAM"})
    client.post("/api/risks", json={"program_id": sibling["id"], "description": "WRONG PROGRAM"})
    owner = client.post("/api/persons", json={"name": "Valence owner", "affiliation": "valence"}).json()
    client.post("/api/commitments", json={"program_id": acct["p"]["id"],
        "description": "RIGHT COMMITMENT", "responsible_party_id": acct["person"]["id"],
        "internal_owner_id": owner["id"], "due_date": _days(5)})
    client.post("/api/commitments", json={"program_id": sibling["id"],
        "description": "WRONG COMMITMENT", "responsible_party_id": acct["person"]["id"],
        "internal_owner_id": owner["id"], "due_date": _days(5)})
    brief = client.get(
        f"/api/accounts/{aid}/pre-call-brief?program_id={acct['p']['id']}").json()
    assert "RIGHT PROGRAM" in brief["markdown"]
    assert "WRONG PROGRAM" not in brief["markdown"]
    assert "RIGHT COMMITMENT" in brief["markdown"]
    assert "WRONG COMMITMENT" not in brief["markdown"]


def test_business_case_rollups_cannot_leak_internal_whitespace(client, acct):
    aid = acct["a"]["id"]
    internal = client.post("/api/whitespace-cells", json={
        "account_id": aid, "segment_id": acct["seg"]["id"], "use_case_id": acct["uc"]["id"],
        "estimated_seats": 5999, "next_action": "SECRET PROCUREMENT TACTIC"}).json()
    client.post(f"/api/whitespace-cells/{internal['id']}/set-fact", json={
        "fact": "evidence_state", "value": "anecdotal", "reason": "internal hypothesis"})
    shared = client.post("/api/whitespace-cells", json={
        "account_id": aid, "segment_id": acct["seg2"]["id"], "use_case_id": acct["uc"]["id"],
        "client_visible": True, "source_reference_id": acct["source"]["id"]}).json()
    client.post(f"/api/whitespace-cells/{shared['id']}/set-fact", json={
        "fact": "evidence_state", "value": "anecdotal", "reason": "shareable pilot"})

    bc = client.get(f"/api/accounts/{aid}/business-case").json()
    assert bc["penetration"]["addressable_seats"] == 4000
    assert bc["ask"]["unpenetrated_seats"] == 4000
    assert {r["population"] for r in bc["ask"]["top_populations"]} == {"Nordics"}
    assert "SECRET PROCUREMENT TACTIC" not in bc["markdown"]


def test_stale_metric_is_not_stamped_current(client, acct):
    aid = acct["a"]["id"]
    d = client.post("/api/metric-definitions", json={
        "name": "Return rate", "stale_after_days": 30}).json()
    old = _days(-90)
    client.post("/api/value-targets", json={
        "account_id": aid, "definition_id": d["id"], "segment_id": acct["seg"]["id"],
        "target_value": .5, "timeframe_end": _days(20), "client_accepted": True,
        "accepted_by_person_id": acct["person"]["id"], "accepted_on": _today(),
        "client_visible": True, "source_reference_id": acct["source"]["id"]})
    client.post("/api/metric-observations", json={
        "definition_id": d["id"], "program_id": acct["p"]["id"],
        "population_segment_id": acct["seg"]["id"], "value": .7,
        "current_through": old, "source_reference_id": acct["source"]["id"]})
    bc = client.get(f"/api/accounts/{aid}/business-case").json()
    assert bc["stamp"]["data_current_through"] == old
    assert any("stale" in gap for gap in bc["stamp"]["missing_or_stale_sources"])
    assert bc["scorecard"][0]["realization"]["value"] is None


def test_draft_is_editable_then_frozen_and_renders_pdf(client, acct):
    doc = client.post(f"/api/accounts/{acct['a']['id']}/documents", json={
        "kind": "business_case"}).json()
    edited = client.patch(f"/api/documents/{doc['id']}", json={
        "title": "Operator-edited title", "body_markdown": "# Reviewed body"}).json()
    assert edited["title"] == "Operator-edited title"
    pdf = client.get(f"/api/documents/{doc['id']}/pdf")
    assert pdf.status_code == 200 and pdf.content.startswith(b"%PDF-")
    client.post(f"/api/documents/{doc['id']}/status", json={
        "status": "reviewed", "reviewed_by": "operator"})
    assert client.patch(f"/api/documents/{doc['id']}", json={
        "body_markdown": "changed after approval"}).status_code == 409


def test_champion_kit_handoff_is_visible_in_champion_pipeline(client, acct):
    client.post("/api/advocacy-events", json={
        "person_id": acct["person"]["id"], "kind": "advocacy_without_us", "occurred_on": _today()})
    client.post("/api/champion-candidates", json={
        "person_id": acct["person"]["id"], "program_id": acct["p"]["id"], "stage": "arm"})
    doc = client.post(f"/api/accounts/{acct['a']['id']}/documents", json={
        "kind": "champion_kit"}).json()
    before = client.get(f"/api/accounts/{acct['a']['id']}/champion-pipeline").json()
    candidate = before["candidates"][0]
    assert candidate["enablement_kits"][0]["document_id"] == doc["id"]
    assert candidate["enablement_kits"][0]["shared_on"] is None
    bundle = client.get(f"/api/accounts/{acct['a']['id']}/export").json()
    assert bundle["tables"]["generated_document_people"][0]["document_id"] == doc["id"]
    client.post(f"/api/documents/{doc['id']}/status", json={
        "status": "sent", "reviewed_by": "operator"})
    after = client.get(f"/api/accounts/{acct['a']['id']}/champion-pipeline").json()
    assert after["candidates"][0]["last_enablement_on"] == _today()
